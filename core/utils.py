import os
import google.generativeai as genai
import chromadb
from pypdf import PdfReader
from django.conf import settings
import io
import docx
import textract
import uuid
from pptx import Presentation

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = "".join(page.extract_text() for page in reader.pages if page.extract_text())
    return text

def extract_text_from_docx(file_path):
    """Extract text from DOCX files using python-docx"""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except Exception as e:
        # Fallback to textract if python-docx fails
        return extract_text_with_textract(file_path)

def extract_text_from_txt(file_path):
    """Extract text from plain text files"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encodings if utf-8 fails
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception:
            return extract_text_with_textract(file_path)

def extract_text_from_pptx(file_path):
    """Extract text from PPTX files using python-pptx"""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        # Fallback to textract if python-pptx fails or for .ppt files
        return extract_text_with_textract(file_path)

def extract_text_with_textract(file_path):
    """Extract text using textract as a fallback method"""
    try:
        return textract.process(file_path).decode('utf-8')
    except Exception as e:
        raise Exception(f"Failed to extract text with textract: {str(e)}")

def extract_text_from_document(file_path):
    """Extract text from a document based on its file extension"""
    _, file_extension = os.path.splitext(file_path.lower())
    
    if file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension == '.docx':
        return extract_text_from_docx(file_path)
    elif file_extension == '.doc' or file_extension == '.rtf':
        return extract_text_with_textract(file_path)
    elif file_extension == '.txt':
        return extract_text_from_txt(file_path)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(file_path)
    elif file_extension == '.ppt': # Older .ppt format, try textract
        return extract_text_with_textract(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {file_extension}")

def chunk_text(text, chunk_size=1000, overlap=100):
    chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size - overlap)]
    return chunks

def get_embeddings(text_chunks):
    result = genai.embed_content(
        model="models/embedding-001",
        content=text_chunks,
        task_type="retrieval_document"
    )
    return result['embedding']

def initialize_chromadb():
    client = chromadb.PersistentClient(path=settings.CHROMA_DB_PATH)
    return client.get_or_create_collection(name=settings.CHROMA_COLLECTION_NAME)

def query_chromadb(collection, query_text, n_results=5):
    query_embedding = genai.embed_content(
        model="models/embedding-001",
        content=query_text,
        task_type="retrieval_query"
    )['embedding']
    results = collection.query(query_embeddings=[query_embedding], n_results=n_results, include=['documents'])
    return results.get('documents', [[]])[0]

def generate_response(query_text, context_chunks):
    context = "\n\n".join(context_chunks)
    prompt = f"""Based ONLY on the following context:\n\n{context}\n\nAnswer the following question:\n{query_text}\n\nAnswer:"""
    response = genai.GenerativeModel('gemini-1.5-flash').generate_content(prompt)
    return response.text if response.text else "No answer available."

def handle_pdf_upload(uploaded_file):
    temp_path = os.path.join(settings.BASE_DIR, 'temp_uploads')
    os.makedirs(temp_path, exist_ok=True)
    # Generate a unique temporary filename to avoid collisions and potential security issues
    original_extension = os.path.splitext(uploaded_file.name)[1]
    temp_filename = f"{uuid.uuid4()}{original_extension}"
    file_path = os.path.join(temp_path, temp_filename)

    try:
        with open(file_path, 'wb') as f:
            f.write(uploaded_file.read())
        return file_path
    except Exception as e:
        raise Exception(f"Failed to save uploaded file: {str(e)}") # Generic error message

